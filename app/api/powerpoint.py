from flask import Blueprint, request, jsonify, send_file
import os
import re
import tempfile
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

powerpoint_bp = Blueprint('powerpoint', __name__)

# Enhanced treatment mappings - French dental terms to actions
TREATMENT_MAPPINGS = {
    # Color changes (crown, onlay, veneer)
    'color_treatments': {
        'cc': 'Couronne céramique',
        'cci': 'Couronne sur implant',
        'couronne': 'Couronne céramique',
        'couronne ceramique': 'Couronne céramique',
        'couronne sur implant': 'Couronne sur implant',
        'onlay': 'Onlay',
        'facette': 'Facette céramique',
        'facette ceramique': 'Facette céramique',
        'veneer': 'Facette céramique',
        'f': 'Facette céramique',
        'cpr': 'Couronne céramique',
        'o': 'Onlay',
        'ceram': 'Facette céramique',
        'ceramique': 'Facette céramique'
    },
    # Icon treatments
    'icon_treatments': {
        'dem': 'Dévitalisation',
        'dém': 'Dévitalisation',
        'devitalisation': 'Dévitalisation',
        'dévitalisation': 'Dévitalisation',
        'tr': 'Traitement endodontique',
        'traitement radiculaire': 'Traitement endodontique',
        'traitement endodontique': 'Traitement endodontique',
        'endo': 'Traitement endodontique',
        'endodontie': 'Traitement endodontique',
        'tenons': 'Tenons',
        'tenon': 'Tenons',
        'ma': 'Moignon adhésif',
        'moignon adhesif': 'Moignon adhésif',
        'moignon adhésif': 'Moignon adhésif',
        'extraction': 'Extraction',
        'ext': 'Extraction',
        'av': 'Extraction',
        'avulsion': 'Extraction',
        'abl': 'Extraction',
        'ablation': 'Extraction',
        'implant': 'Pose d\'implant',
        'imp': 'Pose d\'implant',
        'pose i': 'Pose d\'implant',
        'pose d\'implant': 'Pose d\'implant',
        'pose implant': 'Pose d\'implant',
        'curtage': 'Curetage',
        'curetage': 'Curetage',
        'seance': 'Séance',
        'séance': 'Séance',
        'bnv': 'Blanchissement interne',
        'blanchiment interne': 'Blanchissement interne',
        'blanchissement interne': 'Blanchissement interne',
        'blanch': 'Blanchissement interne',
        'blanchiment': 'Blanchissement interne',
        'gbr': 'Greffe osseuse',
        'greffe osseuse': 'Greffe osseuse',
        'gc': 'Greffe gingivale',
        'greffe gingivale': 'Greffe gingivale',
        'sl': 'Sinus lift',
        'sinus lift': 'Sinus lift',
        'det': 'Détartrage',
        'hd': 'Détartrage',
        'détartrage': 'Détartrage',
        'detartrage': 'Détartrage',
        'sf': 'Scellement de fissure',
        'scellement de fissure': 'Scellement de fissure',
        'dds': 'Dent de sagesse',
        'dent de sagesse': 'Dent de sagesse',
        'fil de cont': 'Fil de contention',
        'fil de contention': 'Fil de contention',
        'dem cc': 'Démonter couronne',
        'dém cc': 'Démonter couronne',
        'te': 'Taille empreinte',
        'taille empreinte': 'Taille empreinte',
        'sc': 'Scellement',
        'scellement': 'Scellement',
        'empr': 'Empreinte',
        'empreinte': 'Empreinte',
        'post-op': 'Post opératoire',
        'post op': 'Post opératoire',
        'prov': 'Provisoire',
        'provisoire': 'Provisoire',
        'm': 'Composite mésial',
        'mesial': 'Composite mésial',
        'mésial': 'Composite mésial',
        'd': 'Composite distal',
        'distal': 'Composite distal',
        'mo': 'Composite mésio-occlusal',
        'do': 'Composite occluso-distal',
        'mod': 'Composite mésio-occluso-distal',
        'l': 'Composite lingual',
        'lingual': 'Composite lingual',
        'p': 'Composite palatin',
        'palatin': 'Composite palatin',
        'v': 'Composite vestibulaire',
        'vestibulaire': 'Composite vestibulaire'
    }
}

# Treatment colors
TREATMENT_COLORS = {
    'Couronne céramique': RGBColor(255, 215, 0),  # Gold
    'Couronne sur implant': RGBColor(255, 165, 0),  # Orange
    'Onlay': RGBColor(192, 192, 192),  # Silver
    'Facette céramique': RGBColor(0, 123, 255)  # Blue
}

def parse_tooth_range(tooth_str):
    """Parse tooth ranges like '12-22' or '11 à 22'"""
    tooth_str = tooth_str.strip()
    
    # Handle ranges with dash: 12-22
    if '-' in tooth_str:
        start, end = tooth_str.split('-')
        start, end = int(start.strip()), int(end.strip())
        return list(range(start, end + 1))
    
    # Handle ranges with 'à': 11 à 22
    elif ' à ' in tooth_str:
        start, end = tooth_str.split(' à ')
        start, end = int(start.strip()), int(end.strip())
        return list(range(start, end + 1))
    
    # Single tooth
    else:
        return [int(tooth_str)]

def enhanced_parse_treatment_text(text):
    """Enhanced parsing with better regex"""
    results = []
    
    # Clean the text
    text = text.lower().strip()
    
    # Enhanced regex patterns
    patterns = [
        # Complex pattern: Plan de TT 11 AV + implant + CC; 22 Implant + CC
        r'plan\s+de\s+t+\s+([^;]+)',
        # Simple pattern: 26 dém. CC + dém. tenons + TR
        r'(\d+(?:\s*[-à]\s*\d+)?)\s*[:\s]*([^;]+)',
        # Alternative pattern: Pour la 26: treatments
        r'pour\s+la\s+(\d+(?:\s*[-à]\s*\d+)?)\s*[:\s]*([^;]+)',
    ]
    
    for pattern in patterns:
        matches = re.findall(pattern, text)
        for match in matches:
            if len(match) == 2:
                tooth_part, treatment_part = match
                
                # Parse tooth numbers (handle ranges)
                try:
                    tooth_numbers = parse_tooth_range(tooth_part)
                except ValueError:
                    continue
                
                # Parse treatments
                treatments = [t.strip() for t in re.split(r'[+&,]', treatment_part) if t.strip()]
                
                for tooth_num in tooth_numbers:
                    for treatment in treatments:
                        treatment = treatment.strip()
                        if treatment:
                            # Determine treatment type and normalize name
                            normalized_treatment = None
                            treatment_type = 'icon'  # default
                            
                            # Check color treatments first
                            for key, value in TREATMENT_MAPPINGS['color_treatments'].items():
                                if key in treatment:
                                    normalized_treatment = value
                                    treatment_type = 'color'
                                    break
                            
                            # Check icon treatments if not found in color treatments
                            if not normalized_treatment:
                                for key, value in TREATMENT_MAPPINGS['icon_treatments'].items():
                                    if key in treatment:
                                        normalized_treatment = value
                                        treatment_type = 'icon'
                                        break
                            
                            # If no mapping found, use original treatment
                            if not normalized_treatment:
                                normalized_treatment = treatment.title()
                            
                            results.append({
                                'tooth': str(tooth_num),
                                'treatment': normalized_treatment,
                                'type': treatment_type,
                                'original': treatment
                            })
    
    return results

def is_valid_tooth_number(tooth_number):
    """Validate tooth number according to FDI system"""
    try:
        tooth_num = int(tooth_number)
        # Valid FDI tooth numbers: 11-18, 21-28, 31-38, 41-48
        valid_ranges = [
            (11, 18), (21, 28), (31, 38), (41, 48)
        ]
        
        for start, end in valid_ranges:
            if start <= tooth_num <= end:
                return True
        return False
    except ValueError:
        return False

def find_tooth_element(slide, tooth_number):
    """Find the tooth element in the slide by searching for tooth name patterns"""
    target_names = [
        f"tooth_{tooth_number}",
        f"Tooth_{tooth_number}",
        f"background_tooth_{tooth_number}",
        f"Background_tooth_{tooth_number}",
        f"tooth{tooth_number}",
        f"Tooth{tooth_number}",
        f"dent_{tooth_number}",
        f"Dent_{tooth_number}"
    ]
    
    def search_in_shapes(shapes):
        # First, check direct shapes by name (faster and more accurate)
        for shape in shapes:
            if hasattr(shape, 'name') and shape.name in target_names:
                return shape
        
        # Then check grouped shapes
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                result = search_in_shapes(shape.shapes)
                if result:
                    return result
        
        # Fallback: check if tooth number appears in shape text
        for shape in shapes:
            if hasattr(shape, 'text') and str(tooth_number) in shape.text:
                return shape
        
        return None
    
    return search_in_shapes(slide.shapes)

def apply_color_treatment(slide, tooth_number, treatment):
    """Apply color treatment to a tooth with enhanced error handling"""
    tooth_element = find_tooth_element(slide, tooth_number)
    
    if not tooth_element:
        return False
    
    try:
        # Get the color for this treatment
        color = TREATMENT_COLORS.get(treatment, RGBColor(128, 128, 128))  # Default gray
        
        # Try multiple approaches to set the color
        # Approach 1: Direct fill
        if hasattr(tooth_element, 'fill'):
            try:
                tooth_element.fill.solid()
                tooth_element.fill.fore_color.rgb = color
                return True
            except Exception:
                pass
        
        # Approach 2: Try line color if fill doesn't work
        if hasattr(tooth_element, 'line'):
            try:
                tooth_element.line.color.rgb = color
                tooth_element.line.width = Pt(3)
                return True
            except Exception:
                pass
        
        # Approach 3: Try text color if it's a text shape
        if hasattr(tooth_element, 'text_frame'):
            try:
                if tooth_element.text_frame.paragraphs:
                    tooth_element.text_frame.paragraphs[0].font.color.rgb = color
                    return True
            except Exception:
                pass
        
        # Approach 4: For grouped shapes, try to color all sub-shapes
        if hasattr(tooth_element, 'shapes'):
            try:
                colored_count = 0
                for sub_shape in tooth_element.shapes:
                    try:
                        if hasattr(sub_shape, 'fill'):
                            sub_shape.fill.solid()
                            sub_shape.fill.fore_color.rgb = color
                            colored_count += 1
                    except:
                        pass
                if colored_count > 0:
                    return True
            except Exception:
                pass
        
        return False
        
    except Exception as e:
        print(f"Error applying color to tooth {tooth_number}: {e}")
        return False

def apply_multiple_icon_treatments(slide, tooth_number, treatments):
    """Apply multiple icon treatments to a tooth with smart positioning"""
    try:
        tooth_element = find_tooth_element(slide, tooth_number)
        if not tooth_element:
            return [False] * len(treatments)
        
        # Get tooth element position and dimensions
        tooth_left = tooth_element.left
        tooth_top = tooth_element.top
        tooth_width = tooth_element.width
        tooth_height = tooth_element.height
        
        # Calculate positions for multiple icons
        icon_size = Inches(0.25)  # Slightly larger for better visibility
        positions = []
        
        if len(treatments) == 1:
            # Single icon - center
            center_x = tooth_left + tooth_width / 2 - icon_size / 2
            center_y = tooth_top + tooth_height / 2 - icon_size / 2
            positions.append((center_x, center_y))
        elif len(treatments) == 2:
            # Two icons - side by side
            spacing = icon_size * 0.1
            total_width = icon_size * 2 + spacing
            start_x = tooth_left + (tooth_width - total_width) / 2
            center_y = tooth_top + tooth_height / 2 - icon_size / 2
            positions.append((start_x, center_y))
            positions.append((start_x + icon_size + spacing, center_y))
        else:
            # Multiple icons - grid layout
            icons_per_row = 2
            rows = (len(treatments) + icons_per_row - 1) // icons_per_row
            
            row_height = icon_size * 0.8
            total_height = row_height * rows
            start_y = tooth_top + (tooth_height - total_height) / 2
            
            for i, treatment in enumerate(treatments):
                row = i // icons_per_row
                col = i % icons_per_row
                
                if row == rows - 1 and len(treatments) % icons_per_row != 0:
                    # Last row with fewer icons - center them
                    remaining_icons = len(treatments) % icons_per_row
                    col_spacing = tooth_width / (remaining_icons + 1)
                    x = tooth_left + col_spacing * (col + 1) - icon_size / 2
                else:
                    col_spacing = tooth_width / (icons_per_row + 1)
                    x = tooth_left + col_spacing * (col + 1) - icon_size / 2
                
                y = start_y + row * row_height
                positions.append((x, y))
        
        # Apply each treatment
        results = []
        for i, treatment in enumerate(treatments):
            try:
                x, y = positions[i]
                success = False
                
                # Try to find and use an icon file first
                icon_path = get_icon_path(treatment)
                
                if icon_path and os.path.exists(icon_path):
                    try:
                        print(f"Adding icon for {treatment}: {icon_path}")
                        # Add the icon image
                        pic = slide.shapes.add_picture(
                            icon_path,
                            x,
                            y,
                            width=icon_size,
                            height=icon_size
                        )
                        success = True
                        print(f"✅ Successfully added icon for {treatment}")
                    except Exception as e:
                        print(f"❌ Error adding icon image for {treatment}: {e}")
                        # Fall back to text if image fails
                        pass
                
                if not success:
                    print(f"Using text fallback for {treatment}")
                    # Fallback to text
                    text_box = slide.shapes.add_textbox(x, y, icon_size, icon_size)
                    
                    # Treatment symbols mapping
                    symbols = {
                        'Extraction': 'EX',
                        'Couronne': 'C',
                        'Implant': 'I',
                        'Pose d\'implant': 'I',
                        'Obturation': 'O',
                        'Endodontie': 'E',
                        'Traitement endodontique': 'TR',
                        'Dévitalisation': 'DÉM',
                        'Prothèse': 'P',
                        'Bridge': 'B',
                        'Facette': 'F',
                        'Détartrage': 'DET',
                        'Blanchiment': 'BL',
                        'Blanchissement interne': 'BNV',
                        'Gingivectomie': 'G',
                        'Greffe osseuse': 'GBR',
                        'Sinus lift': 'SL',
                        'Chirurgie gingivale': 'CG',
                        'Greffe gingivale': 'GC',
                        'Restauration composite': 'RC',
                        'Moignon adhésif': 'MA',
                        'Tenons': 'T',
                        'Curetage': 'C',
                        'Séance': 'S',
                        'Démonter couronne': 'DC',
                        'Taille empreinte': 'TE',
                        'Scellement': 'SC',
                        'Empreinte': 'E',
                        'Post opératoire': 'PO',
                        'Dent de sagesse': 'DS',
                        'Fil de contention': 'FC',
                        'Scellement de fissure': 'SF',
                        'Provisoire': 'P',
                        'Composite mésial': 'M',
                        'Composite distal': 'D',
                        'Composite mésio-occlusal': 'MO',
                        'Composite occluso-distal': 'OD',
                        'Composite mésio-occluso-distal': 'MOD',
                        'Composite lingual': 'L',
                        'Composite palatin': 'P',
                        'Composite vestibulaire': 'V'
                    }
                    
                    text_frame = text_box.text_frame
                    text_frame.text = symbols.get(treatment, treatment[:3].upper())
                    text_frame.margin_left = 0
                    text_frame.margin_right = 0
                    text_frame.margin_top = 0
                    text_frame.margin_bottom = 0
                    
                    # Style the text (smaller for multiple icons)
                    paragraph = text_frame.paragraphs[0]
                    font_size = Pt(10) if len(treatments) == 1 else Pt(8)
                    paragraph.font.size = font_size
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red
                    
                    # Center the text
                    paragraph.alignment = PP_ALIGN.CENTER
                    
                    success = True
                
                results.append(success)
                
            except Exception as e:
                print(f"Icon treatment error for tooth {tooth_number}, treatment {treatment}: {e}")
                results.append(False)
        
        return results
        
    except Exception as e:
        print(f"Multiple icon treatment error for tooth {tooth_number}: {e}")
        return [False] * len(treatments)

def get_icon_path(treatment):
    """Get the path to the icon file for a treatment"""
    # Get the project root directory
    project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    
    # Treatment to icon filename mapping
    icon_mappings = {
        # Main treatments
        'Blanchissement interne': 'bnv.png',
        'Traitement endodontique': 'tr.png',
        'Dévitalisation': 'tr.png',  # Use TR icon for dévitalisation
        'Extraction': 'extraction.png',
        'Pose d\'implant': 'implant.png',
        'Implant': 'implant.png',
        'Greffe osseuse': 'gbr.png',
        'Greffe gingivale': 'gc.png',
        'Moignon adhésif': 'ma.png',
        'Détartrage': 'det.png',
        
        # Additional mappings for common treatments
        'Endodontie': 'tr.png',
        'Avulsion': 'extraction.png',
        'Curetage': 'gc.png',  # Use GC icon for curetage
        'Séance': 'det.png',   # Use DET icon for general séance
        
        # Composite treatments
        'Restauration composite': 'ma.png',  # Use MA as fallback
        'Composite mésial': 'ma.png',
        'Composite distal': 'ma.png',
        'Composite mésio-occlusal': 'ma.png',
        'Composite occluso-distal': 'ma.png',
        'Composite mésio-occluso-distal': 'ma.png',
        'Composite lingual': 'ma.png',
        'Composite palatin': 'ma.png',
        'Composite vestibulaire': 'ma.png',
        
        # Other treatments
        'Tenons': 'ma.png',
        'Chirurgie gingivale': 'gc.png',
        'Sinus lift': 'gbr.png',  # Use GBR for sinus lift
        'Greffe d\'os': 'gbr.png',
        'Greffe de gencive': 'gc.png',
    }
    
    # First try direct mapping
    if treatment in icon_mappings:
        icon_path = os.path.join(project_root, 'static', 'icons', icon_mappings[treatment])
        if os.path.exists(icon_path):
            return icon_path
    
    # Try to create a safe filename from treatment name
    safe_name = re.sub(r'[^\w\s-]', '', treatment.lower())
    safe_name = re.sub(r'[-\s]+', '_', safe_name)
    
    # Common icon file extensions
    extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp']
    
    for ext in extensions:
        icon_path = os.path.join(project_root, 'static', 'icons', f"{safe_name}{ext}")
        if os.path.exists(icon_path):
            return icon_path
    
    return None

def process_powerpoint_treatments(treatments):
    """Process the PowerPoint with the given treatments"""
    try:
        # Get the project root directory
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        
        # Load the PowerPoint template
        template_path = os.path.join(project_root, 'plan.pptx')
        if not os.path.exists(template_path):
            return None, "Template PowerPoint file not found"
        
        prs = Presentation(template_path)
        
        results = []
        
        # Apply treatments to the first slide (assuming dental chart is on first slide)
        if prs.slides:
            slide = prs.slides[0]
            
            # Group treatments by tooth and type
            tooth_treatments = {}
            
            for treatment in treatments:
                tooth = treatment['tooth']
                
                # First check if it's a valid tooth number
                if not is_valid_tooth_number(tooth):
                    results.append({
                        'tooth': tooth,
                        'treatment': treatment['treatment'],
                        'success': False,
                        'error': f"Numéro de dent invalide ({tooth} n'existe pas dans le système FDI)"
                    })
                    continue
                
                if tooth not in tooth_treatments:
                    tooth_treatments[tooth] = {'color': [], 'icon': []}
                
                tooth_treatments[tooth][treatment['type']].append(treatment)
            
            # Process each tooth's treatments
            for tooth, treatments_by_type in tooth_treatments.items():
                print(f"Processing tooth {tooth}...")
                
                # Apply color treatments first (they don't stack)
                for color_treatment in treatments_by_type['color']:
                    success = apply_color_treatment(slide, tooth, color_treatment['treatment'])
                    results.append({
                        'tooth': tooth,
                        'treatment': color_treatment['treatment'],
                        'success': success,
                        'error': f"Élément tooth_{tooth} non trouvé dans le PowerPoint" if not success else None
                    })
                
                # Apply all icon treatments for this tooth at once (with smart positioning)
                if treatments_by_type['icon']:
                    icon_treatments = [t['treatment'] for t in treatments_by_type['icon']]
                    success_list = apply_multiple_icon_treatments(slide, tooth, icon_treatments)
                    
                    for i, (icon_treatment, success) in enumerate(zip(treatments_by_type['icon'], success_list)):
                        results.append({
                            'tooth': tooth,
                            'treatment': icon_treatment['treatment'],
                            'success': success,
                            'error': f"Impossible d'ajouter l'icône sur la dent {tooth}" if not success else None
                        })
        
        # Save the modified presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"plan_modified_{timestamp}.pptx"
        
        # Create temp directory if it doesn't exist
        temp_dir = os.path.join(project_root, 'temp')
        os.makedirs(temp_dir, exist_ok=True)
        output_path = os.path.join(temp_dir, output_filename)
        
        prs.save(output_path)
        
        return output_filename, results
        
    except Exception as e:
        print(f"PowerPoint processing error: {e}")
        return None, str(e)

@powerpoint_bp.route('/generate', methods=['POST'])
def generate_powerpoint():
    """Process PowerPoint generation request"""
    try:
        data = request.json
        text = data.get('treatment_text', '').strip()
        
        if not text:
            return jsonify({'status': 'error', 'message': 'Aucun texte fourni'}), 400
        
        # Parse the treatment text
        treatments = enhanced_parse_treatment_text(text)
        
        if not treatments:
            return jsonify({'status': 'error', 'message': 'Aucun traitement reconnu dans le texte'}), 400
        
        # Process the PowerPoint
        output_file, results = process_powerpoint_treatments(treatments)
        
        if output_file:
            return jsonify({
                'status': 'success',
                'treatments': results,
                'file_path': output_file
            })
        else:
            return jsonify({'status': 'error', 'message': results}), 500
            
    except Exception as e:
        print(f"Error in generate_powerpoint: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500

@powerpoint_bp.route('/download/<path:filename>', methods=['GET'])
def download_powerpoint(filename):
    """Download generated PowerPoint file"""
    try:
        # Secure filename
        filename = os.path.basename(filename)
        
        # Get the project root directory (parent of app directory)
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
        file_path = os.path.join(project_root, 'temp', filename)
        
        # Debug logging
        print(f"Attempting to download file: {filename}")
        print(f"Project root: {project_root}")
        print(f"Full file path: {file_path}")
        print(f"File exists: {os.path.exists(file_path)}")
        
        if os.path.exists(file_path):
            return send_file(
                file_path,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=filename
            )
        else:
            return jsonify({
                'status': 'error',
                'message': f'Fichier non trouvé: {filename}'
            }), 404
            
    except Exception as e:
        print(f"Error downloading PowerPoint: {e}")
        return jsonify({
            'status': 'error',
            'message': str(e)
        }), 500