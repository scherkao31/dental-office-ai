import os
import tempfile
from datetime import datetime
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PowerPointService:
    """Service for generating PowerPoint presentations"""
    
    def generate_treatment_presentation(self, treatment_data: Dict, dental_schema: Optional[Dict] = None) -> str:
        """Generate PowerPoint presentation for treatment plan"""
        prs = Presentation()
        
        # Title slide
        self._add_title_slide(prs, treatment_data['patient_name'])
        
        # Treatment overview slide
        self._add_overview_slide(prs, treatment_data)
        
        # Dental schema slide if provided
        if dental_schema:
            self._add_dental_schema_slide(prs, dental_schema)
        
        # Treatment sequence slides
        for sequence in treatment_data.get('sequences', []):
            self._add_sequence_slide(prs, sequence)
        
        # Cost summary slide
        self._add_cost_summary_slide(prs, treatment_data)
        
        # Save presentation
        filename = f"treatment_plan_{treatment_data['patient_name'].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
        filepath = os.path.join(tempfile.gettempdir(), filename)
        prs.save(filepath)
        
        return filepath
    
    def _add_title_slide(self, prs: Presentation, patient_name: str):
        """Add title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Plan de Traitement Dentaire"
        subtitle.text = f"Patient: {patient_name}\n{datetime.now().strftime('%d %B %Y')}"
    
    def _add_overview_slide(self, prs: Presentation, treatment_data: Dict):
        """Add treatment overview slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Vue d'ensemble du traitement"
        
        # Add content
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        
        # Summary points
        p = tf.add_paragraph()
        p.text = f"Nombre de séquences: {len(treatment_data.get('sequences', []))}"
        p.level = 0
        
        total_treatments = sum(len(seq.get('treatments', [])) for seq in treatment_data.get('sequences', []))
        p = tf.add_paragraph()
        p.text = f"Nombre total de traitements: {total_treatments}"
        p.level = 0
        
        if 'duration' in treatment_data:
            p = tf.add_paragraph()
            p.text = f"Durée estimée: {treatment_data['duration']}"
            p.level = 0
        
        if 'total_cost' in treatment_data:
            p = tf.add_paragraph()
            p.text = f"Coût total estimé: {treatment_data['total_cost']} CHF"
            p.level = 0
    
    def _add_dental_schema_slide(self, prs: Presentation, dental_schema: Dict):
        """Add dental schema visualization slide"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.75)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "Schéma Dentaire"
        
        # Format title
        for paragraph in title_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(24)
                run.font.bold = True
        
        # Add dental schema content (simplified representation)
        content_top = Inches(1.5)
        content_box = slide.shapes.add_textbox(left, content_top, width, Inches(5))
        content_frame = content_box.text_frame
        
        # Upper teeth
        p = content_frame.add_paragraph()
        p.text = "Maxillaire (Supérieur)"
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        
        # Add tooth numbers and conditions
        if 'upper_teeth' in dental_schema:
            p = content_frame.add_paragraph()
            p.text = ' '.join([f"{t['number']}" for t in dental_schema['upper_teeth']])
            p.alignment = PP_ALIGN.CENTER
        
        # Lower teeth
        p = content_frame.add_paragraph()
        p.text = "\nMandibule (Inférieur)"
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        
        if 'lower_teeth' in dental_schema:
            p = content_frame.add_paragraph()
            p.text = ' '.join([f"{t['number']}" for t in dental_schema['lower_teeth']])
            p.alignment = PP_ALIGN.CENTER
    
    def _add_sequence_slide(self, prs: Presentation, sequence: Dict):
        """Add slide for treatment sequence"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = f"Séquence {sequence['sequence']}: {sequence['title']}"
        
        # Add treatments
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        
        for treatment in sequence.get('treatments', []):
            p = tf.add_paragraph()
            p.text = treatment['description']
            p.level = 0
            
            # Add details as sub-bullets
            if 'duration' in treatment:
                p = tf.add_paragraph()
                p.text = f"Durée: {treatment['duration']}"
                p.level = 1
            
            if 'cost' in treatment:
                p = tf.add_paragraph()
                p.text = f"Coût: {treatment['cost']} CHF"
                p.level = 1
            
            if 'notes' in treatment:
                p = tf.add_paragraph()
                p.text = f"Notes: {treatment['notes']}"
                p.level = 1
    
    def _add_cost_summary_slide(self, prs: Presentation, treatment_data: Dict):
        """Add cost summary slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Résumé des Coûts"
        
        # Add cost breakdown
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        
        # Add costs by sequence
        for sequence in treatment_data.get('sequences', []):
            sequence_cost = sum(t.get('cost', 0) for t in sequence.get('treatments', []))
            p = tf.add_paragraph()
            p.text = f"{sequence['title']}: {sequence_cost} CHF"
            p.level = 0
        
        # Add total
        tf.add_paragraph()  # Empty line
        p = tf.add_paragraph()
        p.text = f"Coût Total: {treatment_data.get('total_cost', 0)} CHF"
        p.level = 0
        
        # Make total bold
        for run in p.runs:
            run.font.bold = True
            run.font.size = Pt(16)
        
        # Add payment options if available
        if 'payment_options' in treatment_data:
            tf.add_paragraph()  # Empty line
            p = tf.add_paragraph()
            p.text = "Options de paiement disponibles"
            p.level = 0
            
            for option in treatment_data['payment_options']:
                p = tf.add_paragraph()
                p.text = option
                p.level = 1
    
    def create_dental_schema_presentation(self, teeth_data: List[Dict], patient_name: str) -> str:
        """Create a presentation focused on dental schema"""
        prs = Presentation()
        
        # Title slide
        self._add_title_slide(prs, patient_name)
        
        # Dental condition overview
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "État Dentaire Actuel"
        
        # Group teeth by condition
        conditions = {}
        for tooth in teeth_data:
            condition = tooth.get('condition', 'healthy')
            if condition not in conditions:
                conditions[condition] = []
            conditions[condition].append(tooth['number'])
        
        # Add conditions to slide
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        
        for condition, teeth in conditions.items():
            p = tf.add_paragraph()
            p.text = f"{condition.title()}: {', '.join(map(str, teeth))}"
            p.level = 0
        
        # Save presentation
        filename = f"dental_schema_{patient_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
        filepath = os.path.join(tempfile.gettempdir(), filename)
        prs.save(filepath)
        
        return filepath